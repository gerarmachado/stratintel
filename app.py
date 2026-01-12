import streamlit as st
import google.generativeai as genai
from langchain_google_genai import ChatGoogleGenerativeAI
import pypdf
from docx import Document
from fpdf import FPDF
from pptx import Presentation # Para PowerPoint
from pptx.util import Inches, Pt
from io import BytesIO
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
import yt_dlp
import os
import time
import datetime
# Agente de Búsqueda
from langchain_community.tools import DuckDuckGoSearchRun

# --- CONFIGURACIÓN DE PÁGINA ---
st.set_page_config(page_title="StratIntel V12 (Commander)", page_icon="🎖️", layout="wide")

# ==========================================
# 🔐 SISTEMA DE LOGIN
# ==========================================
def check_password():
    """Retorna `True` si el usuario tiene la contraseña correcta."""
    def password_entered():
        if st.session_state["username"] in st.secrets["passwords"] and \
           st.session_state["password"] == st.secrets["passwords"][st.session_state["username"]]:
            st.session_state["password_correct"] = True
            del st.session_state["password"]
        else:
            st.session_state["password_correct"] = False

    if st.session_state.get("password_correct", False):
        return True

    st.markdown("## 🎖️ StratIntel: Acceso Restringido")
    st.text_input("Usuario", key="username")
    st.text_input("Contraseña", type="password", on_change=password_entered, key="password")
    
    if "password_correct" in st.session_state and not st.session_state["password_correct"]:
        st.error("❌ Credenciales inválidas")
    return False

if not check_password():
    st.stop()

# ==========================================
# ⚙️ CONFIGURACIÓN Y MODELO
# ==========================================
API_KEY_FIJA = "" 
if "GOOGLE_API_KEY" in st.secrets:
    API_KEY_FIJA = st.secrets["GOOGLE_API_KEY"]

MODELO_ACTUAL = "gemini-3-flash-preview"  

# ==========================================
# 🧠 BASE DE DATOS DE CONOCIMIENTO
# ==========================================
DB_CONOCIMIENTO = {
    "✨ RECOMENDACIÓN AUTOMÁTICA": {
        "desc": "La IA decide la mejor estrategia.",
        "preguntas": ["Hallazgos críticos.", "Evaluación de riesgos.", "Resumen Ejecutivo (BLUF).", "Patrones ocultos."]
    },
    "Niveles de Análisis (Barry Buzan)": {
        "desc": "Seguridad Multisectorial (Militar, Política, Económica, Societal, Ambiental).",
        "preguntas": ["Nivel Sistémico (Polaridad).", "Nivel Estatal (Presiones internas).", "Nivel Individual (Líderes).", "Seguridad Societal (Identidad)."]
    },
    "Evolución de la Cooperación (Axelrod)": {
        "desc": "Teoría de Juegos.",
        "preguntas": ["Sombra del Futuro.", "Reciprocidad (Tit-for-Tat).", "Detección de Trampas.", "Estructura de Pagos."]
    },
    "Análisis FODA (SWOT) Intel": {
        "desc": "Enfoque de Inteligencia.",
        "preguntas": ["Vulnerabilidades (Debilidades).", "Amenazas inminentes.", "Estrategia Maxi-Mini.", "Fortalezas vs Oportunidades."]
    },
    "Análisis Geopolítico (PMESII-PT)": {
        "desc": "Variables del entorno operativo.",
        "preguntas": ["Político-Militar.", "Infraestructura/Social.", "Desglose PMESII-PT."]
    },
    "Análisis DIME (Poder Nacional)": {
        "desc": "Diplomático, Informacional, Militar, Económico.",
        "preguntas": ["Poder Económico.", "Aislamiento Diplomático.", "Guerra de Info.", "Disuasión Militar."]
    },
    "Análisis de Hipótesis (ACH)": {
        "desc": "Validación de Hipótesis.",
        "preguntas": ["Matriz de Hipótesis.", "Evidencia diagnóstica.", "Intelligence Gaps.", "Decepción."]
    },
    "Abogado del Diablo": {
        "desc": "Pensamiento crítico.",
        "preguntas": ["Desafío a la conclusión.", "Defensa del actor 'irracional'."]
    },
    "Escenarios Prospectivos": {
        "desc": "Cono de Plausibilidad.",
        "preguntas": ["4 Escenarios (Mejor/Peor/Híbrido/Wild).", "Drivers clave."]
    },
    "Centro de Gravedad (COG)": {
        "desc": "Clausewitz.",
        "preguntas": ["COG Estratégico.", "Capacidades Críticas.", "Vulnerabilidades Críticas."]
    },
    "Matriz CARVER": {
        "desc": "Selección de Objetivos.",
        "preguntas": ["Criticidad/Vulnerabilidad.", "Efecto sistémico.", "Recuperabilidad."]
    }
}

# --- GESTIÓN DE ESTADO ---
if 'api_key' not in st.session_state: st.session_state['api_key'] = ""
if 'texto_analisis' not in st.session_state: st.session_state['texto_analisis'] = ""
if 'origen_dato' not in st.session_state: st.session_state['origen_dato'] = "Ninguno"

# --- FUNCIONES DE PROCESAMIENTO ---

def buscar_en_web(query):
    """Agente de búsqueda usando DuckDuckGo (Gratis)"""
    try:
        search = DuckDuckGoSearchRun()
        resultados = search.run(query)
        return resultados
    except Exception as e:
        return f"Error en búsqueda web: {e}"

def procesar_archivos_pdf(archivos):
    texto_total = ""
    nombres = []
    for archivo in archivos:
        reader = pypdf.PdfReader(archivo)
        texto_pdf = "".join([p.extract_text() for p in reader.pages])
        texto_total += f"\n--- ARCHIVO: {archivo.name} ---\n{texto_pdf}\n"
        nombres.append(archivo.name)
    return texto_total, str(nombres)

def procesar_archivos_docx(archivos):
    texto_total = ""
    nombres = []
    for archivo in archivos:
        doc = Document(archivo)
        texto_doc = "\n".join([para.text for para in doc.paragraphs])
        texto_total += f"\n--- ARCHIVO: {archivo.name} ---\n{texto_doc}\n"
        nombres.append(archivo.name)
    return texto_total, str(nombres)

def obtener_texto_web(url):
    try:
        h = {'User-Agent': 'Mozilla/5.0'}
        r = requests.get(url, headers=h, timeout=15)
        s = BeautifulSoup(r.content, 'html.parser')
        for script in s(["script", "style"]): script.extract()
        return s.get_text(separator='\n')
    except Exception as e: return f"Error: {e}"

def procesar_youtube(url, api_key):
    vid = url.split("v=")[-1].split("&")[0] if "v=" in url else url.split("/")[-1]
    try:
        t = YouTubeTranscriptApi.get_transcript(vid, languages=['es', 'en'])
        return " ".join([i['text'] for i in t]), "Subtítulos"
    except:
        st.info(f"Modo Multimodal (Audio)...")
        opts = {'format': 'bestaudio/best', 'outtmpl': '%(id)s.%(ext)s', 'postprocessors': [{'key': 'FFmpegExtractAudio','preferredcodec': 'mp3'}], 'quiet': True}
        try:
            with yt_dlp.YoutubeDL(opts) as ydl:
                info = ydl.extract_info(url, download=True)
                fname = f"{info['id']}.mp3"
            genai.configure(api_key=api_key)
            myfile = genai.upload_file(fname)
            while myfile.state.name == "PROCESSING": time.sleep(2); myfile = genai.get_file(myfile.name)
            model = genai.GenerativeModel(MODELO_ACTUAL)
            res = model.generate_content([myfile, "Transcribe audio."])
            if os.path.exists(fname): os.remove(fname)
            myfile.delete()
            return res.text, "Audio IA"
        except Exception as e: return f"Error: {e}", "Error"

# --- FUNCIONES DE REPORTE (PPTX, PDF, DOCX) ---

def limpiar_texto(t):
    if not t: return ""
    reps = {"✨": "", "🚀": "", "⚠️": "[!]", "✅": "[OK]", "🛡️": "", "🔒": "", "🎖️": ""}
    for k,v in reps.items(): t = t.replace(k,v)
    return t.encode('latin-1', 'replace').decode('latin-1')

def crear_pptx(texto, tecnicas, fuente):
    prs = Presentation()
    
    # Diapositiva de Título
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Informe de Inteligencia StratIntel V12"
    subtitle.text = f"Fuente: {fuente}\nGenerado por IA"

    # Procesar texto para diapositivas
    # Dividimos por técnica (asumiendo que usamos encabezados ##)
    secciones = texto.split("## 📌")
    
    for seccion in secciones:
        if not seccion.strip(): continue
        
        lineas = seccion.strip().split("\n")
        titulo_seccion = lineas[0].strip() # El nombre de la técnica
        contenido = "\n".join(lineas[1:])[:1000] # Limitamos texto por slide
        
        # Crear Slide
        layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(layout)
        
        # Título
        title = slide.shapes.title
        title.text = titulo_seccion
        
        # Contenido (Bullet points simples)
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = contenido

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

class PDFReport(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 12)
        self.cell(0, 10, 'StratIntel Report V12', 0, 1, 'C')
        self.ln(5)
    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 7)
        self.cell(0, 10, 'Generado por IA. Uso Confidencial.', 0, 0, 'C')

def crear_pdf(texto, tecnicas, fuente):
    pdf = PDFReport()
    pdf.add_page()
    pdf.set_font("Arial", "B", 10)
    pdf.multi_cell(0, 5, limpiar_texto(f"Fuente: {fuente}\nTécnicas: {tecnicas}"))
    pdf.ln(5)
    pdf.set_font("Arial", "", 10)
    pdf.multi_cell(0, 5, limpiar_texto(texto))
    return pdf.output(dest='S').encode('latin-1', 'replace')

def crear_word(texto, tecnicas, fuente):
    doc = Document()
    doc.add_heading('StratIntel Intelligence Report V12', 0)
    doc.add_paragraph(f"Fuente: {fuente}").bold = True
    doc.add_paragraph(f"Técnicas Aplicadas: {tecnicas}").bold = True
    for l in texto.split('\n'):
        if l.startswith('#'): doc.add_heading(l.replace('#','').strip(), level=2)
        else: doc.add_paragraph(l)
    
    aviso = doc.add_paragraph()
    aviso_runner = aviso.add_run(
        "\n\n------------------------------------------------------------\n"
        "AVISO DE RESPONSABILIDAD: Documento generado por IA. "
        "Verificar hallazgos críticos."
    )
    aviso_runner.font.size = 8
    aviso_runner.italic = True
    
    b = BytesIO(); doc.save(b); b.seek(0)
    return b

# --- INTERFAZ ---
st.sidebar.title("🎖️ StratIntel V12")
st.sidebar.caption("Commander Edition | Multi-Select")
st.sidebar.markdown("---")

if API_KEY_FIJA:
    st.session_state['api_key'] = API_KEY_FIJA
    genai.configure(api_key=API_KEY_FIJA)
    st.sidebar.success(f"✅ Conectado ({MODELO_ACTUAL})")
else:
    if not st.session_state['api_key']:
        k = st.sidebar.text_input("🔑 API KEY:", type="password")
        if k: st.session_state['api_key'] = k; genai.configure(api_key=k); st.rerun()

# --- SELECTOR MULTI-TÉCNICA (NUEVO) ---
st.sidebar.subheader("🎯 Configuración de Misión")
tecnicas_seleccionadas = st.sidebar.multiselect(
    "Selecciona Técnicas (Máx. 3):",
    options=list(DB_CONOCIMIENTO.keys()),
    max_selections=3,
    help="Combina metodologías compatibles para un análisis coherente."
)

temp = st.sidebar.slider("Creatividad", 0.0, 1.0, 0.4)

if st.sidebar.button("🔒 Salir"):
    del st.session_state["password_correct"]
    st.rerun()

st.title("🎖️ StratIntel | División de Análisis")
st.markdown("**Sistema de Apoyo a la Decisión (DSS) v12.0**")

# --- CARGA ---
t1, t2, t3, t4, t5 = st.tabs(["📂 Multi-PDF", "📝 Multi-DOCX", "🌐 Web", "📺 YouTube", "✍️ Manual"])

with t1:
    pdfs = st.file_uploader("Subir PDFs", type="pdf", accept_multiple_files=True)
    if pdfs and st.button("Procesar PDF"):
        txt, n = procesar_archivos_pdf(pdfs)
        st.session_state['texto_analisis'] = txt
        st.session_state['origen_dato'] = f"PDFs: {n}"
        st.success(f"✅ {len(pdfs)} archivos.")

with t2:
    docs = st.file_uploader("Subir Words", type="docx", accept_multiple_files=True)
    if docs and st.button("Procesar DOCX"):
        txt, n = procesar_archivos_docx(docs)
        st.session_state['texto_analisis'] = txt
        st.session_state['origen_dato'] = f"DOCXs: {n}"
        st.success(f"✅ {len(docs)} archivos.")

with t3:
    u = st.text_input("URL Noticia")
    if st.button("Extraer"): st.session_state['texto_analisis'] = obtener_texto_web(u); st.session_state['origen_dato'] = f"Web: {u}"; st.success("Cargado")

with t4:
    y = st.text_input("URL YouTube")
    if st.button("Video"):
        with st.spinner("Procesando..."):
            txt, m = procesar_youtube(y, st.session_state['api_key'])
            if m!="Error": st.session_state['texto_analisis']=txt; st.session_state['origen_dato']=f"YT: {y}"; st.success("Cargado")
            else: st.error(txt)

with t5:
    m = st.text_area("Manual")
    if st.button("Fijar"): st.session_state['texto_analisis']=m; st.session_state['origen_dato']="Manual"; st.success("Cargado")

st.markdown("---")
if st.session_state['texto_analisis']:
    st.info(f"📂 Fuente Activa: **{st.session_state['origen_dato']}**")
    with st.expander("Ver Datos"): st.write(st.session_state['texto_analisis'][:1000] + "...")

# --- EJECUCIÓN V12 ---
st.header("Generación de Inteligencia")

if not st.session_state['api_key'] or not st.session_state['texto_analisis']:
    st.warning("⚠️ Carga datos para comenzar.")
else:
    c1, c2 = st.columns([1, 2])
    
    with c1:
        st.markdown("### ⚙️ Parámetros")
        if not tecnicas_seleccionadas:
            st.info("👈 Selecciona al menos 1 técnica en la barra lateral.")
        
        # CHECKBOX PARA AGENTE DE BÚSQUEDA
        usar_internet = st.checkbox("🌐 Activar Búsqueda Web en Vivo", help="La IA buscará información actualizada en DuckDuckGo antes de analizar.")
        
        pir = st.text_area("Requerimiento Específico (PIR):", placeholder="Opcional: Define un enfoque...", height=100)

    with c2:
        if st.button("🚀 EJECUTAR MISIÓN (MULTI-TECNICA)", type="primary", use_container_width=True, disabled=len(tecnicas_seleccionadas)==0):
            try:
                genai.configure(api_key=st.session_state['api_key'])
                model = genai.GenerativeModel(MODELO_ACTUAL)
                ctx = st.session_state['texto_analisis']
                
                # FASE 1: BÚSQUEDA WEB (SI ESTÁ ACTIVADA)
                contexto_web = ""
                if usar_internet:
                    with st.status("🌐 Agente buscando en internet...", expanded=True) as status:
                        query_search = f"{pir} {st.session_state['origen_dato']}" if pir else f"Análisis estratégico sobre {st.session_state['origen_dato']}"
                        res_web = buscar_en_web(query_search)
                        contexto_web = f"\n\n--- 🌐 INFORMACIÓN ACTUALIZADA DE INTERNET ---\n{res_web}\n------------------------------------------\n"
                        status.update(label="✅ Búsqueda completada", state="complete", expanded=False)
                
                # FASE 2: BUCLE DE ANÁLISIS
                informe_final = f"# INFORME DE INTELIGENCIA V12\nFECHA: {datetime.datetime.now().strftime('%d/%m/%Y')}\nFUENTE: {st.session_state['origen_dato']}\n\n"
                
                progreso = st.progress(0)
                
                for i, tec in enumerate(tecnicas_seleccionadas):
                    st.caption(f"Analizando: {tec}...")
                    
                    # PROMPT CONSTRUIDO CON DATOS + WEB
                    prompt = f"""
                    ACTÚA COMO: Especialista en Derecho y Política Internacional y Analista Senior.
                    TAREA: Realizar un análisis profundo usando la metodología: {tec}.
                    PIR (Requerimiento): {pir}
                    
                    INSTRUCCIONES:
                    1. Profundidad académica (2-3 párrafos por punto).
                    2. Cita fuentes.
                    3. Integra la información web si es relevante.
                    
                    CONTEXTO DOCUMENTAL:
                    {ctx}
                    
                    {contexto_web}
                    """
                    
                    res = model.generate_content(prompt, generation_config=genai.types.GenerationConfig(temperature=temp))
                    informe_final += f"\n\n## 📌 {tec}\n{res.text}\n\n---\n"
                    
                    progreso.progress((i + 1) / len(tecnicas_seleccionadas))
                    time.sleep(2) # Pausa técnica
                
                st.session_state['res'] = informe_final
                st.session_state['tecnicas_usadas'] = ", ".join(tecnicas_seleccionadas)
                st.success("✅ Misión Completada")
                st.markdown(informe_final)

            except Exception as e: st.error(f"Error: {e}")

# DESCARGAS
if 'res' in st.session_state:
    st.markdown("---")
    st.subheader("📥 Descargar Paquete de Inteligencia")
    col1, col2, col3 = st.columns(3)
    
    # WORD
    doc_b = crear_word(st.session_state['res'], st.session_state['tecnicas_usadas'], st.session_state['origen_dato'])
    col1.download_button("📄 Informe WORD", doc_b, "Informe_V12.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    
    # PDF
    try:
        pdf_b = crear_pdf(st.session_state['res'], st.session_state['tecnicas_usadas'], st.session_state['origen_dato'])
        col2.download_button("📑 Informe PDF", bytes(pdf_b), "Informe_V12.pdf", "application/pdf")
    except: pass
    
    # POWERPOINT (NUEVO)
    try:
        pptx_b = crear_pptx(st.session_state['res'], st.session_state['tecnicas_usadas'], st.session_state['origen_dato'])
        col3.download_button("📊 Presentación PPTX", pptx_b, "Presentacion_V12.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e: col3.error(f"Error PPTX: {e}")
